#!/usr/bin/env python3
"""
測試 PDF 生成功能
"""

import sys
from pathlib import Path
from datetime import datetime

def test_markdown():
    """測試 Markdown 模組"""
    try:
        import markdown
        md = "# 測試\n這是測試文件"
        html = markdown.markdown(md)
        print("✅ Markdown 模組正常")
        return True
    except ImportError:
        print("❌ Markdown 模組未安裝")
        return False

def test_reportlab():
    """測試 ReportLab (PDF 生成)"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        
        output_dir = Path.home() / ".openclaw" / "workspace" / "test-output"
        output_dir.mkdir(exist_ok=True)
        
        pdf_file = output_dir / "test-report.pdf"
        
        c = canvas.Canvas(str(pdf_file), pagesize=A4)
        c.setFont("Helvetica", 24)
        c.drawString(100, 700, "OpenClaw PDF Test")
        c.setFont("Helvetica", 12)
        c.drawString(100, 680, f"Generated at: {datetime.now().isoformat()}")
        c.drawString(100, 660, "PDF generation is working!")
        c.save()
        
        print(f"✅ ReportLab PDF 生成成功：{pdf_file}")
        return True
    except ImportError:
        print("❌ ReportLab 模組未安裝")
        return False
    except Exception as e:
        print(f"❌ PDF 生成失敗：{e}")
        return False

def test_openpyxl():
    """測試 OpenPyXL (Excel 生成)"""
    try:
        from openpyxl import Workbook
        
        wb = Workbook()
        ws = wb.active
        ws.title = "任務統計"
        
        ws['A1'] = "任務名稱"
        ws['B1'] = "狀態"
        ws['C1'] = "Agent"
        
        ws['A2'] = "測試任務 1"
        ws['B2'] = "已完成"
        ws['C2'] = "main"
        
        ws['A3'] = "測試任務 2"
        ws['B3'] = "執行中"
        ws['C3'] = "main"
        
        output_dir = Path.home() / ".openclaw" / "workspace" / "test-output"
        output_dir.mkdir(exist_ok=True)
        
        excel_file = output_dir / "test-tasks.xlsx"
        wb.save(str(excel_file))
        
        print(f"✅ OpenPyXL Excel 生成成功：{excel_file}")
        return True
    except ImportError:
        print("❌ OpenPyXL 模組未安裝")
        return False

def test_pptx():
    """測試 Python-pptx (PowerPoint 生成)"""
    try:
        from pptx import Presentation
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "OpenClaw 任務報告"
        subtitle.text = f"生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        output_dir = Path.home() / ".openclaw" / "workspace" / "test-output"
        output_dir.mkdir(exist_ok=True)
        
        pptx_file = output_dir / "test-report.pptx"
        prs.save(str(pptx_file))
        
        print(f"✅ Python-pptx PowerPoint 生成成功：{pptx_file}")
        return True
    except ImportError:
        print("❌ Python-pptx 模組未安裝")
        return False

def test_pandoc():
    """測試 Pandoc"""
    import subprocess
    try:
        result = subprocess.run(['pandoc', '--version'], capture_output=True, text=True)
        if result.returncode == 0:
            version = result.stdout.split('\n')[0]
            print(f"✅ Pandoc: {version}")
            return True
    except:
        print("❌ Pandoc 未安裝")
        return False
    return False

def test_wkhtmltopdf():
    """測試 wkhtmltopdf"""
    import subprocess
    try:
        result = subprocess.run(['wkhtmltopdf', '--version'], capture_output=True, text=True)
        if result.returncode == 0:
            version = result.stdout.strip()
            print(f"✅ wkhtmltopdf: {version}")
            return True
    except:
        print("❌ wkhtmltopdf 未安裝")
        return False
    return False

def test_sqlite3():
    """測試 SQLite3"""
    import subprocess
    try:
        result = subprocess.run(['sqlite3', '--version'], capture_output=True, text=True)
        if result.returncode == 0:
            version = result.stdout.strip()
            print(f"✅ SQLite3: {version}")
            return True
    except:
        print("❌ SQLite3 未安裝")
        return False
    return False

def test_jq():
    """測試 jq"""
    import subprocess
    try:
        result = subprocess.run(['jq', '--version'], capture_output=True, text=True)
        if result.returncode == 0:
            version = result.stdout.strip()
            print(f"✅ jq: {version}")
            return True
    except:
        print("❌ jq 未安裝")
        return False
    return False

def main():
    print("🦞 OpenClaw PDF/Office 生成測試")
    print("=" * 50)
    print()
    
    results = []
    
    # 測試系統工具
    print("📦 系統工具測試:")
    results.append(test_pandoc())
    results.append(test_wkhtmltopdf())
    results.append(test_sqlite3())
    results.append(test_jq())
    print()
    
    # 測試 Python 模組
    print("🐍 Python 模組測試:")
    results.append(test_markdown())
    results.append(test_reportlab())
    results.append(test_openpyxl())
    results.append(test_pptx())
    print()
    
    # 總結
    total = len(results)
    passed = sum(results)
    
    print("=" * 50)
    print(f"測試結果：{passed}/{total} 通過")
    print(f"完成度：{passed * 100 // total}%")
    
    if passed == total:
        print("\n🎉 所有測試通過！系統已就緒！")
    else:
        print(f"\n⚠️  有 {total - passed} 個測試未通過，請執行安裝腳本。")
    
    print()
    print("測試檔案位置：")
    print(f"  ~/.openclaw/workspace/test-output/")

if __name__ == "__main__":
    main()
